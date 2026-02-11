"""
Slide Converter - 画像→編集可能PowerPoint変換ツール（プロトタイプ）
Kirigami風のAI画像解析 + python-pptx生成
"""

import os
import json
import base64
import re
import uuid
import traceback
from io import BytesIO
from pathlib import Path

from flask import Flask, request, jsonify, send_file, render_template_string
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max

UPLOAD_DIR = Path("/home/claude/slide-converter/uploads")
OUTPUT_DIR = Path("/home/claude/slide-converter/outputs")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ============================================================
# Claude API呼び出し（Vision）
# ============================================================

def analyze_slide_image(image_path: str, api_key: str) -> dict:
    """Claude API Visionで画像を解析し、要素をJSON形式で取得"""
    import urllib.request

    # 画像をbase64エンコード
    with open(image_path, "rb") as f:
        image_data = base64.standard_b64encode(f.read()).decode("utf-8")

    # 拡張子からメディアタイプ判定
    ext = Path(image_path).suffix.lower()
    media_types = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".gif": "image/gif", ".webp": "image/webp"}
    media_type = media_types.get(ext, "image/png")

    prompt = """この画像はプレゼンテーションスライドです。画像内の全要素を解析し、以下のJSON形式で正確に返してください。
JSONのみを返し、他のテキストやマークダウンのコードブロックは含めないでください。

{
  "slide_width_px": <画像の横幅ピクセル推定>,
  "slide_height_px": <画像の縦幅ピクセル推定>,
  "background": {
    "type": "solid" or "gradient",
    "color": "#RRGGBB",
    "gradient_end": "#RRGGBB"
  },
  "elements": [
    {
      "type": "text",
      "content": "テキスト内容",
      "x_percent": 0-100,
      "y_percent": 0-100,
      "width_percent": 0-100,
      "height_percent": 0-100,
      "font_size_pt": 数値,
      "font_color": "#RRGGBB",
      "bold": true/false,
      "italic": true/false,
      "align": "left"/"center"/"right",
      "background_color": "#RRGGBB" or null
    },
    {
      "type": "shape",
      "shape_type": "rectangle"/"rounded_rectangle"/"oval"/"line",
      "x_percent": 0-100,
      "y_percent": 0-100,
      "width_percent": 0-100,
      "height_percent": 0-100,
      "fill_color": "#RRGGBB",
      "border_color": "#RRGGBB" or null,
      "border_width": 数値 or null
    },
    {
      "type": "image_region",
      "description": "画像の説明（グラフ、写真など）",
      "x_percent": 0-100,
      "y_percent": 0-100,
      "width_percent": 0-100,
      "height_percent": 0-100
    }
  ]
}

重要なルール:
- 座標はすべてスライド全体に対するパーセンテージ（0-100）で指定
- テキストは正確に読み取り、潰れた文字は文脈から推測して修復
- 背景の上に重なっている要素を全て検出
- 要素はz-order（背面→前面）の順に並べる
- 図形の中にテキストがある場合、図形とテキストを別要素として出力
- font_size_ptは画像内での見た目から推定（タイトル: 28-44pt, 本文: 14-18pt, キャプション: 10-12pt）
- JSONのみ出力。説明文やコードブロック記号は不要"""

    payload = {
        "model": "claude-sonnet-4-5-20250514",
        "max_tokens": 4096,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": media_type,
                            "data": image_data
                        }
                    },
                    {
                        "type": "text",
                        "text": prompt
                    }
                ]
            }
        ]
    }

    headers = {
        "Content-Type": "application/json",
        "x-api-key": api_key,
        "anthropic-version": "2023-06-01"
    }

    req = urllib.request.Request(
        "https://api.anthropic.com/v1/messages",
        data=json.dumps(payload).encode("utf-8"),
        headers=headers,
        method="POST"
    )

    with urllib.request.urlopen(req, timeout=120) as resp:
        result = json.loads(resp.read().decode("utf-8"))

    # レスポンスからテキスト部分を取得
    text_content = ""
    for block in result.get("content", []):
        if block.get("type") == "text":
            text_content += block["text"]

    # JSONをパース（コードブロックの場合も対応）
    text_content = text_content.strip()
    text_content = re.sub(r'^```json\s*', '', text_content)
    text_content = re.sub(r'\s*```$', '', text_content)

    return json.loads(text_content)


# ============================================================
# PowerPoint生成
# ============================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """#RRGGBB → RGBColor"""
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def build_pptx(analysis: dict, original_image_path: str = None) -> BytesIO:
    """解析結果からPowerPointファイルを生成"""
    prs = Presentation()

    # 16:9レイアウト
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 空白レイアウトを使用
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    # 背景色設定
    bg = analysis.get("background", {})
    bg_color = bg.get("color", "#FFFFFF")
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(bg_color)

    # 各要素を配置
    elements = analysis.get("elements", [])

    for elem in elements:
        try:
            elem_type = elem.get("type", "")
            x = int(slide_w * elem.get("x_percent", 0) / 100)
            y = int(slide_h * elem.get("y_percent", 0) / 100)
            w = int(slide_w * elem.get("width_percent", 10) / 100)
            h = int(slide_h * elem.get("height_percent", 10) / 100)

            # 最小サイズ保証
            w = max(w, Emu(100000))
            h = max(h, Emu(100000))

            if elem_type == "text":
                txBox = slide.shapes.add_textbox(x, y, w, h)
                tf = txBox.text_frame
                tf.word_wrap = True

                # 背景色がある場合
                bg_col = elem.get("background_color")
                if bg_col:
                    txBox_fill = txBox.fill
                    txBox_fill.solid()
                    txBox_fill.fore_color.rgb = hex_to_rgb(bg_col)

                p = tf.paragraphs[0]
                p.text = elem.get("content", "")

                # フォント設定
                font_size = elem.get("font_size_pt", 16)
                font_color = elem.get("font_color", "#333333")
                bold = elem.get("bold", False)
                italic = elem.get("italic", False)
                align = elem.get("align", "left")

                run = p.runs[0] if p.runs else p.add_run()
                if not p.runs:
                    run.text = elem.get("content", "")
                    p.text = ""

                run.font.size = Pt(font_size)
                run.font.color.rgb = hex_to_rgb(font_color)
                run.font.bold = bold
                run.font.italic = italic

                align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
                p.alignment = align_map.get(align, PP_ALIGN.LEFT)

            elif elem_type == "shape":
                shape_type = elem.get("shape_type", "rectangle")
                shape_map = {
                    "rectangle": MSO_SHAPE.RECTANGLE,
                    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
                    "oval": MSO_SHAPE.OVAL,
                    "line": MSO_SHAPE.RECTANGLE,  # lineはthin rectで代用
                }
                mso_shape = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)

                if shape_type == "line":
                    h = max(Emu(30000), h)  # 線は薄く

                shape = slide.shapes.add_shape(mso_shape, x, y, w, h)

                fill_color = elem.get("fill_color")
                if fill_color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
                else:
                    shape.fill.background()

                border_color = elem.get("border_color")
                border_width = elem.get("border_width")
                if border_color:
                    shape.line.color.rgb = hex_to_rgb(border_color)
                    shape.line.width = Pt(border_width or 1)
                else:
                    shape.line.fill.background()

            elif elem_type == "image_region":
                # 画像領域はプレースホルダーとして矩形+テキストで表現
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
                shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                shape.line.width = Pt(1)

                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                desc = elem.get("description", "画像")
                p.text = f"[{desc}]"
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0]
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
                run.font.italic = True

        except Exception as e:
            print(f"要素配置エラー: {elem_type} - {e}")
            traceback.print_exc()
            continue

    # ファイルをメモリに書き出し
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ============================================================
# デモモード（API不要）
# ============================================================

def demo_analysis() -> dict:
    """APIなしで動作確認するためのサンプルデータ"""
    return {
        "slide_width_px": 1280,
        "slide_height_px": 720,
        "background": {"type": "solid", "color": "#1E2761"},
        "elements": [
            {
                "type": "shape",
                "shape_type": "rectangle",
                "x_percent": 0, "y_percent": 0,
                "width_percent": 100, "height_percent": 100,
                "fill_color": "#1E2761",
                "border_color": None
            },
            {
                "type": "text",
                "content": "売上報告 2025年度",
                "x_percent": 10, "y_percent": 8,
                "width_percent": 80, "height_percent": 12,
                "font_size_pt": 40,
                "font_color": "#FFFFFF",
                "bold": True, "italic": False,
                "align": "center",
                "background_color": None
            },
            {
                "type": "shape",
                "shape_type": "rounded_rectangle",
                "x_percent": 5, "y_percent": 25,
                "width_percent": 42, "height_percent": 55,
                "fill_color": "#2A3A8F",
                "border_color": "#4A5ABF",
                "border_width": 2
            },
            {
                "type": "text",
                "content": "Q1: ¥12,500,000\nQ2: ¥15,800,000\nQ3: ¥18,200,000\nQ4: ¥22,100,000",
                "x_percent": 8, "y_percent": 30,
                "width_percent": 36, "height_percent": 45,
                "font_size_pt": 18,
                "font_color": "#CADCFC",
                "bold": False, "italic": False,
                "align": "left",
                "background_color": None
            },
            {
                "type": "image_region",
                "description": "売上推移 棒グラフ",
                "x_percent": 53, "y_percent": 25,
                "width_percent": 42, "height_percent": 55
            },
            {
                "type": "text",
                "content": "© 2025 Sample Corp. All rights reserved.",
                "x_percent": 10, "y_percent": 88,
                "width_percent": 80, "height_percent": 8,
                "font_size_pt": 10,
                "font_color": "#7788AA",
                "bold": False, "italic": True,
                "align": "center",
                "background_color": None
            }
        ]
    }


# ============================================================
# HTMLテンプレート
# ============================================================

HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Slide Converter - 画像→PowerPoint変換</title>
<style>
  @import url('https://fonts.googleapis.com/css2?family=Noto+Sans+JP:wght@300;400;600;700&family=JetBrains+Mono:wght@400;600&display=swap');

  :root {
    --bg-primary: #0a0e1a;
    --bg-secondary: #111827;
    --bg-card: #1a2236;
    --border: #2a3654;
    --text-primary: #e2e8f0;
    --text-secondary: #8896b3;
    --accent: #6366f1;
    --accent-glow: rgba(99, 102, 241, 0.3);
    --success: #10b981;
    --warning: #f59e0b;
    --error: #ef4444;
  }

  * { margin: 0; padding: 0; box-sizing: border-box; }

  body {
    font-family: 'Noto Sans JP', sans-serif;
    background: var(--bg-primary);
    color: var(--text-primary);
    min-height: 100vh;
    overflow-x: hidden;
  }

  /* Background grid effect */
  body::before {
    content: '';
    position: fixed;
    top: 0; left: 0; right: 0; bottom: 0;
    background-image:
      linear-gradient(rgba(99, 102, 241, 0.03) 1px, transparent 1px),
      linear-gradient(90deg, rgba(99, 102, 241, 0.03) 1px, transparent 1px);
    background-size: 40px 40px;
    pointer-events: none;
    z-index: 0;
  }

  .container {
    max-width: 860px;
    margin: 0 auto;
    padding: 2rem 1.5rem;
    position: relative;
    z-index: 1;
  }

  /* Header */
  .header {
    text-align: center;
    margin-bottom: 2.5rem;
  }

  .header h1 {
    font-size: 2rem;
    font-weight: 700;
    letter-spacing: -0.02em;
    background: linear-gradient(135deg, #818cf8, #6366f1, #a78bfa);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    margin-bottom: 0.5rem;
  }

  .header p {
    color: var(--text-secondary);
    font-size: 0.9rem;
    font-weight: 300;
  }

  .header .tag {
    display: inline-block;
    margin-top: 0.6rem;
    padding: 0.2rem 0.7rem;
    background: rgba(99, 102, 241, 0.15);
    border: 1px solid rgba(99, 102, 241, 0.3);
    border-radius: 999px;
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.7rem;
    color: #818cf8;
  }

  /* Card */
  .card {
    background: var(--bg-card);
    border: 1px solid var(--border);
    border-radius: 12px;
    padding: 1.8rem;
    margin-bottom: 1.5rem;
  }

  .card-title {
    font-size: 0.8rem;
    font-weight: 600;
    text-transform: uppercase;
    letter-spacing: 0.08em;
    color: var(--text-secondary);
    margin-bottom: 1rem;
  }

  /* API Key */
  .api-key-row {
    display: flex;
    gap: 0.8rem;
    align-items: center;
  }

  .api-key-row input {
    flex: 1;
    padding: 0.7rem 1rem;
    background: var(--bg-secondary);
    border: 1px solid var(--border);
    border-radius: 8px;
    color: var(--text-primary);
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.85rem;
    outline: none;
    transition: border-color 0.2s;
  }

  .api-key-row input:focus {
    border-color: var(--accent);
    box-shadow: 0 0 0 3px var(--accent-glow);
  }

  .api-key-row input::placeholder {
    color: #4a5568;
  }

  /* Upload area */
  .upload-area {
    border: 2px dashed var(--border);
    border-radius: 12px;
    padding: 2.5rem 1.5rem;
    text-align: center;
    cursor: pointer;
    transition: all 0.3s;
    position: relative;
  }

  .upload-area:hover, .upload-area.dragover {
    border-color: var(--accent);
    background: rgba(99, 102, 241, 0.05);
  }

  .upload-area .icon {
    font-size: 2.5rem;
    margin-bottom: 0.8rem;
    opacity: 0.6;
  }

  .upload-area p {
    color: var(--text-secondary);
    font-size: 0.9rem;
  }

  .upload-area .formats {
    font-size: 0.75rem;
    color: #4a5568;
    margin-top: 0.4rem;
  }

  .upload-area input[type="file"] {
    position: absolute;
    inset: 0;
    opacity: 0;
    cursor: pointer;
  }

  /* Preview */
  .preview {
    margin-top: 1rem;
    display: none;
  }

  .preview img {
    max-width: 100%;
    max-height: 300px;
    border-radius: 8px;
    border: 1px solid var(--border);
  }

  .preview .filename {
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.8rem;
    color: var(--text-secondary);
    margin-top: 0.5rem;
  }

  /* Buttons */
  .btn {
    display: inline-flex;
    align-items: center;
    gap: 0.5rem;
    padding: 0.75rem 1.5rem;
    border: none;
    border-radius: 8px;
    font-family: 'Noto Sans JP', sans-serif;
    font-size: 0.9rem;
    font-weight: 600;
    cursor: pointer;
    transition: all 0.2s;
  }

  .btn-primary {
    background: var(--accent);
    color: white;
    width: 100%;
    justify-content: center;
    margin-top: 1rem;
  }

  .btn-primary:hover:not(:disabled) {
    background: #5558e6;
    box-shadow: 0 4px 20px var(--accent-glow);
  }

  .btn-primary:disabled {
    opacity: 0.4;
    cursor: not-allowed;
  }

  .btn-demo {
    background: transparent;
    border: 1px solid var(--border);
    color: var(--text-secondary);
    width: 100%;
    justify-content: center;
    margin-top: 0.5rem;
  }

  .btn-demo:hover {
    border-color: var(--text-secondary);
    color: var(--text-primary);
  }

  .btn-download {
    background: var(--success);
    color: white;
    width: 100%;
    justify-content: center;
    margin-top: 1rem;
    text-decoration: none;
    font-family: 'Noto Sans JP', sans-serif;
    font-size: 0.9rem;
    font-weight: 600;
    padding: 0.75rem 1.5rem;
    border-radius: 8px;
    display: none;
  }

  .btn-download:hover {
    background: #059669;
  }

  /* Status */
  .status {
    margin-top: 1rem;
    padding: 0.8rem 1rem;
    border-radius: 8px;
    font-size: 0.85rem;
    display: none;
  }

  .status.processing {
    display: block;
    background: rgba(99, 102, 241, 0.1);
    border: 1px solid rgba(99, 102, 241, 0.3);
    color: #818cf8;
  }

  .status.success {
    display: block;
    background: rgba(16, 185, 129, 0.1);
    border: 1px solid rgba(16, 185, 129, 0.3);
    color: var(--success);
  }

  .status.error {
    display: block;
    background: rgba(239, 68, 68, 0.1);
    border: 1px solid rgba(239, 68, 68, 0.3);
    color: var(--error);
  }

  /* JSON viewer */
  .json-viewer {
    margin-top: 1rem;
    display: none;
  }

  .json-viewer summary {
    cursor: pointer;
    font-size: 0.8rem;
    color: var(--text-secondary);
    margin-bottom: 0.5rem;
  }

  .json-viewer pre {
    background: var(--bg-secondary);
    border: 1px solid var(--border);
    border-radius: 8px;
    padding: 1rem;
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.75rem;
    color: #a5b4cf;
    overflow-x: auto;
    max-height: 400px;
    overflow-y: auto;
  }

  /* Spinner */
  @keyframes spin { to { transform: rotate(360deg); } }
  .spinner {
    display: inline-block;
    width: 16px; height: 16px;
    border: 2px solid rgba(129, 140, 248, 0.3);
    border-top-color: #818cf8;
    border-radius: 50%;
    animation: spin 0.8s linear infinite;
    vertical-align: middle;
    margin-right: 0.5rem;
  }

  /* Footer */
  .footer {
    text-align: center;
    margin-top: 2rem;
    font-size: 0.75rem;
    color: #3a4560;
  }
</style>
</head>
<body>
<div class="container">
  <div class="header">
    <h1>✂️ Slide Converter</h1>
    <p>画像を、編集可能なPowerPointへ</p>
    <span class="tag">prototype v0.1</span>
  </div>

  <!-- API Key -->
  <div class="card">
    <div class="card-title">🔑 Claude API Key</div>
    <div class="api-key-row">
      <input type="password" id="apiKey" placeholder="sk-ant-api03-..." />
    </div>
    <p style="font-size:0.75rem; color:#4a5568; margin-top:0.5rem;">
      APIキーはサーバーに保存されません。画像解析のみに使用されます。
    </p>
  </div>

  <!-- Upload -->
  <div class="card">
    <div class="card-title">📎 スライド画像をアップロード</div>
    <div class="upload-area" id="uploadArea">
      <div class="icon">📄</div>
      <p>クリックまたはドラッグ＆ドロップ</p>
      <p class="formats">PNG / JPG / WEBP（最大16MB）</p>
      <input type="file" id="fileInput" accept="image/png,image/jpeg,image/webp" />
    </div>
    <div class="preview" id="preview">
      <img id="previewImg" src="" alt="preview" />
      <div class="filename" id="filename"></div>
    </div>
  </div>

  <!-- Convert -->
  <div class="card">
    <div class="card-title">⚡ 変換</div>
    <button class="btn btn-primary" id="convertBtn" disabled onclick="convert(false)">
      PowerPointに変換
    </button>
    <button class="btn btn-demo" onclick="convert(true)">
      デモモード（APIキー不要）
    </button>

    <div class="status" id="status"></div>

    <a class="btn btn-download" id="downloadBtn" href="#" download>
      📥 PowerPointをダウンロード
    </a>

    <div class="json-viewer" id="jsonViewer">
      <details>
        <summary>解析結果JSON</summary>
        <pre id="jsonContent"></pre>
      </details>
    </div>
  </div>

  <div class="footer">
    Slide Converter Prototype — Claude Vision API + python-pptx
  </div>
</div>

<script>
const fileInput = document.getElementById('fileInput');
const uploadArea = document.getElementById('uploadArea');
const preview = document.getElementById('preview');
const previewImg = document.getElementById('previewImg');
const filename = document.getElementById('filename');
const convertBtn = document.getElementById('convertBtn');
const status = document.getElementById('status');
const downloadBtn = document.getElementById('downloadBtn');
const jsonViewer = document.getElementById('jsonViewer');
const jsonContent = document.getElementById('jsonContent');

let selectedFile = null;

// Drag & Drop
uploadArea.addEventListener('dragover', (e) => { e.preventDefault(); uploadArea.classList.add('dragover'); });
uploadArea.addEventListener('dragleave', () => { uploadArea.classList.remove('dragover'); });
uploadArea.addEventListener('drop', (e) => {
  e.preventDefault();
  uploadArea.classList.remove('dragover');
  if (e.dataTransfer.files.length) {
    fileInput.files = e.dataTransfer.files;
    handleFile(e.dataTransfer.files[0]);
  }
});

fileInput.addEventListener('change', (e) => {
  if (e.target.files.length) handleFile(e.target.files[0]);
});

function handleFile(file) {
  selectedFile = file;
  const reader = new FileReader();
  reader.onload = (e) => {
    previewImg.src = e.target.result;
    preview.style.display = 'block';
    filename.textContent = file.name + ' (' + (file.size / 1024).toFixed(1) + ' KB)';
  };
  reader.readAsDataURL(file);
  convertBtn.disabled = false;
  downloadBtn.style.display = 'none';
  jsonViewer.style.display = 'none';
}

async function convert(demo) {
  const apiKey = document.getElementById('apiKey').value.trim();
  if (!demo && !apiKey) {
    showStatus('error', 'APIキーを入力してください');
    return;
  }
  if (!demo && !selectedFile) {
    showStatus('error', '画像をアップロードしてください');
    return;
  }

  showStatus('processing', '<span class="spinner"></span>変換中... Claude APIで画像を解析しています');
  convertBtn.disabled = true;
  downloadBtn.style.display = 'none';
  jsonViewer.style.display = 'none';

  const formData = new FormData();
  if (!demo && selectedFile) formData.append('image', selectedFile);
  formData.append('api_key', apiKey);
  formData.append('demo', demo ? '1' : '0');

  try {
    const resp = await fetch('/convert', { method: 'POST', body: formData });
    const data = await resp.json();

    if (data.success) {
      showStatus('success', '✅ 変換完了！ PowerPointファイルが生成されました');
      downloadBtn.href = '/download/' + data.filename;
      downloadBtn.style.display = 'flex';
      jsonContent.textContent = JSON.stringify(data.analysis, null, 2);
      jsonViewer.style.display = 'block';
    } else {
      showStatus('error', '❌ エラー: ' + data.error);
    }
  } catch (e) {
    showStatus('error', '❌ 通信エラー: ' + e.message);
  }

  convertBtn.disabled = false;
}

function showStatus(type, html) {
  status.className = 'status ' + type;
  status.innerHTML = html;
}
</script>
</body>
</html>
"""


# ============================================================
# Routes
# ============================================================

@app.route("/")
def index():
    return render_template_string(HTML_TEMPLATE)


@app.route("/convert", methods=["POST"])
def convert():
    try:
        demo = request.form.get("demo", "0") == "1"
        api_key = request.form.get("api_key", "").strip()

        if demo:
            # デモモード
            analysis = demo_analysis()
            output = build_pptx(analysis)
        else:
            # 画像アップロード確認
            if "image" not in request.files:
                return jsonify({"success": False, "error": "画像がアップロードされていません"})

            file = request.files["image"]
            if not file.filename:
                return jsonify({"success": False, "error": "ファイルが選択されていません"})

            if not api_key:
                return jsonify({"success": False, "error": "APIキーが入力されていません"})

            # 画像保存
            file_id = str(uuid.uuid4())[:8]
            ext = Path(file.filename).suffix.lower() or ".png"
            image_path = UPLOAD_DIR / f"{file_id}{ext}"
            file.save(str(image_path))

            # Claude APIで解析
            analysis = analyze_slide_image(str(image_path), api_key)

            # PPTX生成
            output = build_pptx(analysis, str(image_path))

        # ファイル保存
        out_id = str(uuid.uuid4())[:8]
        out_filename = f"converted_{out_id}.pptx"
        out_path = OUTPUT_DIR / out_filename
        with open(out_path, "wb") as f:
            f.write(output.read())

        return jsonify({
            "success": True,
            "filename": out_filename,
            "analysis": analysis
        })

    except json.JSONDecodeError as e:
        return jsonify({"success": False, "error": f"AIの応答をパースできませんでした: {str(e)}"})
    except Exception as e:
        traceback.print_exc()
        return jsonify({"success": False, "error": str(e)})


@app.route("/download/<filename>")
def download(filename):
    # パストラバーサル対策
    safe_name = Path(filename).name
    file_path = OUTPUT_DIR / safe_name
    if not file_path.exists():
        return "File not found", 404
    return send_file(
        str(file_path),
        as_attachment=True,
        download_name=safe_name,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


if __name__ == "__main__":
    print("=" * 50)
    print("  Slide Converter - Prototype v0.1")
    print("  http://localhost:8081")
    print("=" * 50)
    app.run(host="0.0.0.0", port=8081, debug=True)
